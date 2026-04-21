# actualizar_pptx.py
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Paths
base_dir = os.path.dirname(os.path.abspath(__file__))
pptx_path = os.path.join(base_dir, "presentacion_tfm.pptx")
output_dir = os.path.join(base_dir, "outputs")

# Abrir la presentación original
prs = Presentation(pptx_path)

# Lista de gráficos a insertar (png + html)
graficos = [
    ("edad_histograma.png", "edad_histograma.html", "Distribución de Edades"),
    ("sexo_pie.png", "sexo_pie.html", "Distribución por Sexo"),
    ("top_regiones.png", "top_regiones.html", "Top 10 Regiones de Pasajeros"),
]

# Insertar cada gráfico
insertados = []
for png, html, titulo in graficos:
    png_path = os.path.join(output_dir, png)
    html_path = os.path.join(output_dir, html)

    if os.path.exists(png_path):
        # Añadir nueva diapositiva con diseño en blanco
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Título
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = titulo
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 51, 102)

        # Imagen
        slide.shapes.add_picture(png_path, Inches(1), Inches(1.5), width=Inches(7))

        # Nota al pie con enlace al HTML
        if os.path.exists(html_path):
            left = Inches(0.5)
            top = Inches(6.5)
            width = Inches(9)
            height = Inches(0.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = f"Ver versión interactiva: {html}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0, 102, 204)

        insertados.append(titulo)

# Guardar sobrescribiendo el original
prs.save(pptx_path)

# Resumen en consola
print("✅ Presentación actualizada correctamente.")
print("Se insertaron los siguientes gráficos:")
for titulo in insertados:
    print(f" - {titulo}")